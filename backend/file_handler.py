"""
文件处理服务
"""
import os
import magic
import logging
from datetime import datetime
from typing import Optional
from PIL import Image
from pathlib import Path

logger = logging.getLogger(__name__)

class FileHandler:
    def __init__(self, upload_folder: str, preview_folder: str):
        self.upload_folder = upload_folder
        self.preview_folder = preview_folder
        self._ensure_directories()
    
    def _ensure_directories(self):
        """确保目录存在"""
        Path(self.upload_folder).mkdir(parents=True, exist_ok=True)
        Path(self.preview_folder).mkdir(parents=True, exist_ok=True)
    
    def allowed_file(self, filename: str, allowed_extensions: set) -> bool:
        """检查文件扩展名是否允许"""
        return '.' in filename and \
            filename.rsplit('.', 1)[1].lower() in allowed_extensions
    
    def get_file_type(self, file_path: str) -> str:
        """获取文件的MIME类型"""
        try:
            mime = magic.Magic(mime=True)
            file_type = mime.from_file(file_path)
            return file_type
        except Exception as e:
            logger.warning(f"无法检测文件类型: {e}")
            return 'application/octet-stream'
    
    def get_file_extension(self, filename: str) -> str:
        """获取文件扩展名"""
        return filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
    
    def format_file_size(self, size_in_bytes: int) -> str:
        """格式化文件大小"""
        for unit in ['B', 'KB', 'MB', 'GB']:
            if size_in_bytes < 1024:
                return f"{size_in_bytes:.2f} {unit}"
            size_in_bytes /= 1024
        return f"{size_in_bytes:.2f} TB"
    
    def generate_preview(
        self,
        file_path: str,
        output_name: str,
        width: int = 800,
        height: int = 1000
    ) -> Optional[str]:
        """
        生成文件预览图
        
        支持: PDF, 图片, 文本, Office文档(通过转换)
        """
        try:
            file_type = self.get_file_type(file_path)
            extension = self.get_file_extension(file_path)
            
            # 图片文件直接生成预览
            if file_type.startswith('image/'):
                return self._generate_image_preview(file_path, output_name, width, height)
            
            # PDF文件
            elif file_type == 'application/pdf' or extension == 'pdf':
                return self._generate_pdf_preview(file_path, output_name, width, height)
            
            # 文本文件
            elif file_type.startswith('text/') or extension == 'txt':
                return self._generate_text_preview(file_path, output_name, width, height)
            
            # Office文档（轻量级预览）
            elif extension in ['doc', 'docx', 'xls', 'xlsx', 'ppt', 'pptx']:
                return self._generate_office_preview(file_path, output_name, width, height)
            
            # 暂时不支持其他格式的预览
            else:
                logger.info(f"文件类型 {file_type} 暂不支持预览生成")
                return None
        
        except Exception as e:
            logger.error(f"生成预览失败: {e}")
            return None
    
    def _generate_office_preview(
        self,
        file_path: str,
        output_name: str,
        width: int,
        height: int
    ) -> Optional[str]:
        """
        生成Office文档轻量级预览（提取文字生成图片）
        """
        try:
            from docx import Document
            import openpyxl
            from pptx import Presentation
            
            extension = self.get_file_extension(file_path)
            text_content = ""
            
            # 提取文字内容
            if extension in ['doc', 'docx']:
                doc = Document(file_path)
                for para in doc.paragraphs[:20]:  # 限制提取前20段
                    text_content += para.text + "\n"
                    
            elif extension in ['xls', 'xlsx']:
                wb = openpyxl.load_workbook(file_path)
                ws = wb.active
                for row in ws.iter_rows(max_row=20, max_col=5):
                    row_text = " | ".join([str(cell.value) for cell in row if cell.value])
                    if row_text.strip():
                        text_content += row_text + "\n"
                        
            elif extension in ['ppt', 'pptx']:
                prs = Presentation(file_path)
                for slide in prs.slides[:5]:  # 限制提取前5页
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_content += shape.text + "\n"
            
            # 生成文字预览图
            if text_content.strip():
                return self._text_to_preview_image(text_content, output_name, width, height)
            
            return None
            
        except ImportError as e:
            logger.warning(f"Office库未安装: {e}")
            return None
        except Exception as e:
            logger.error(f"Office预览生成失败: {e}")
            return None
    
    def _text_to_preview_image(
        self,
        text_content: str,
        output_name: str,
        width: int,
        height: int
    ) -> str:
        """将文字内容转换为预览图"""
        from PIL import Image, ImageDraw, ImageFont
        
        img = Image.new('RGB', (width, height), color='white')
        draw = ImageDraw.Draw(img)
        
        # 使用默认字体（轻量化，不依赖系统字体）
        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 12)
        except:
            font = ImageFont.load_default()
        
        # 限制文字长度
        max_chars = 2000
        text_content = text_content[:max_chars] + ("..." if len(text_content) > max_chars else "")
        
        # 绘制标题
        title = f"文档预览 (前 {len(text_content)} 字符)"
        try:
            title_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 14)
        except:
            title_font = font
        draw.text((20, 20), title, fill='#333', font=title_font)
        
        # 绘制内容
        y = 50
        lines = text_content.split('\n')
        for line in lines:
            if y > height - 40:
                break
            draw.text((20, y), line.strip()[:100], fill='#666', font=font)
            y += 16
        
        output_path = os.path.join(self.preview_folder, f"{output_name}.png")
        img.save(output_path, 'PNG')
        
        return f"/previews/{output_name}.png"
    
    def _generate_image_preview(
        self,
        file_path: str,
        output_name: str,
        width: int,
        height: int
    ) -> str:
        """生成图片预览"""
        with Image.open(file_path) as img:
            # 保持宽高比缩放
            img.thumbnail((width, height))
            
            output_path = os.path.join(self.preview_folder, f"{output_name}.png")
            img.save(output_path, 'PNG')
        
        return f"/previews/{output_name}.png"
    
    def _generate_pdf_preview(self, file_path: str, output_name: str, width: int, height: int) -> str:
        """
        生成PDF预览(需要pdf2image和poppler)
        
        注意: Windows环境下需要安装poppler
        """
        try:
            from pdf2image import convert_from_path
            
            images = convert_from_path(file_path, dpi=72, first_page=1, last_page=1)
            
            if images:
                output_path = os.path.join(self.preview_folder, f"{output_name}.png")
                images[0].thumbnail((width, height))
                images[0].save(output_path, 'PNG')
                return f"/previews/{output_name}.png"
        
        except ImportError:
            logger.warning("pdf2image未安装，无法生成PDF预览")
        except Exception as e:
            logger.error(f"PDF预览生成失败: {e}")
        
        return None
    
    def _generate_text_preview(self, file_path: str, output_name: str, width: int, height: int) -> str:
        """生成文本预览"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read(5000)  # 限制读取前5000字符
            
            # 创建文本预览图
            img = Image.new('RGB', (width, height), color='white')
            from PIL import ImageDraw, ImageFont
            
            draw = ImageDraw.Draw(img)
            
            try:
                font = ImageFont.truetype("arial.ttf", 14)
            except:
                font = ImageFont.load_default()
            
            # 绘制文本
            lines = []
            words = content.split('\n')
            current_line = ""
            
            for word in words:
                test_line = current_line + word + "\n"
                bbox = draw.textbbox((0, 0), test_line, font=font)
                if bbox[2] < width - 40:
                    current_line = test_line
                else:
                    lines.append(current_line)
                    current_line = word + "\n"
                
                if len(lines) > height // 20:
                    break
            
            lines.append(current_line)
            
            y_position = 20
            for line in lines:
                if y_position > height - 40:
                    break
                draw.text((20, y_position), line, fill='black', font=font)
                y_position += 20
            
            output_path = os.path.join(self.preview_folder, f"{output_name}.png")
            img.save(output_path, 'PNG')
            
            return f"/previews/{output_name}.png"
        
        except Exception as e:
            logger.error(f"文本预览生成失败: {e}")
            return None
    
    def convert_to_pdf(self, file_path: str) -> Optional[str]:
        """
        将文件转换为PDF格式
        
        支持: Office文档, 图片, 文本
        """
        try:
            extension = self.get_file_extension(file_path)
            output_path = os.path.splitext(file_path)[0] + '.pdf'
            
            # 图片转PDF
            if extension.lower() in ['jpg', 'jpeg', 'png', 'gif', 'bmp']:
                return self._image_to_pdf(file_path, output_path)
            
            # Office文档转PDF
            elif extension.lower() in ['doc', 'docx', 'xls', 'xlsx', 'ppt', 'pptx']:
                return self._office_to_pdf(file_path, output_path)
            
            # 文本转PDF
            elif extension.lower() == 'txt':
                return self._text_to_pdf(file_path, output_path)
            
            return file_path  # 已经是PDF或其他格式
        
        except Exception as e:
            logger.error(f"文件转换失败: {e}")
            return None
    
    def _image_to_pdf(self, file_path: str, output_path: str) -> str:
        """图片转PDF"""
        images = []
        extensions = ['jpg', 'jpeg', 'png', 'gif', 'bmp']
        
        if '*' in file_path or '?' in file_path:
            from glob import glob
            files = sorted(glob(file_path))
            for f in files:
                if self.get_file_extension(f) in extensions:
                    images.append(Image.open(f))
        else:
            images.append(Image.open(file_path))
        
        if images:
            images[0].save(
                output_path,
                save_all=True,
                append_images=images[1:]
            )
            return output_path
        
        return file_path
    
    def _office_to_pdf(self, file_path: str, output_path: str) -> str:
        """Office文档转PDF"""
        extension = self.get_file_extension(file_path).lower()
        logger.info(f"开始转换Office文档: {file_path} -> {output_path}")
        
        try:
            # DOCX 转 PDF
            if extension in ['doc', 'docx']:
                result = self._docx_to_pdf(file_path, output_path)
                if result != file_path:
                    return result
                    
            # XLSX 转 PDF
            elif extension in ['xls', 'xlsx']:
                result = self._xlsx_to_pdf(file_path, output_path)
                if result != file_path:
                    return result
                    
            # PPTX 转 PDF
            elif extension in ['ppt', 'pptx']:
                result = self._pptx_to_pdf(file_path, output_path)
                if result != file_path:
                    return result
            
            logger.error(f"Office文档转换失败: {file_path}")
            return file_path
            
        except Exception as e:
            logger.error(f"Office文档转换异常: {e}")
            return file_path
    
    def _docx_to_pdf(self, file_path: str, output_path: str) -> str:
        """DOCX转PDF"""
        try:
            from docx import Document
            from weasyprint import HTML
            
            doc = Document(file_path)
            
            # 创建HTML（保留基本格式）
            html_content = "<html><head><meta charset='utf-8'></head><body>"
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    # 检测标题样式
                    if para.style.name.startswith('Heading'):
                        level = int(para.style.name.replace('Heading', '')) or 2
                        html_content += f"<h{level}>{text}</h{level}>"
                    else:
                        html_content += f"<p>{text}</p>"
            html_content += "</body></html>"
            
            HTML(string=html_content).write_pdf(output_path)
            logger.info(f"DOCX转PDF成功: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"DOCX转PDF失败: {e}")
            return file_path
    
    def _xlsx_to_pdf(self, file_path: str, output_path: str) -> str:
        """XLSX转PDF"""
        try:
            import openpyxl
            from openpyxl.utils.dataframe import dataframe_to_rows
            from weasyprint import HTML
            
            wb = openpyxl.load_workbook(file_path)
            ws = wb.active
            
            # 获取所有数据
            data = []
            for row in ws.iter_rows(max_row=ws.max_row or 100, max_col=ws.max_column or 10):
                data.append([cell.value for cell in row])
            
            if not data:
                return file_path
            
            # 创建HTML表格
            html_content = "<html><head><meta charset='utf-8'></head><body>"
            html_content += "<h2>Excel表格</h2>"
            html_content += "<table border='1' style='border-collapse: collapse; width: 100%;'>"
            
            for r_idx, row in enumerate(data):
                html_content += "<tr>"
                for c_idx, value in enumerate(row):
                    tag = "th" if r_idx == 0 else "td"
                    cell_value = str(value or '')
                    html_content += f"<{tag}>{cell_value}</{tag}>"
                html_content += "</tr>"
            
            html_content += "</table></body></html>"
            
            HTML(string=html_content).write_pdf(output_path)
            logger.info(f"XLSX转PDF成功: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"XLSX转PDF失败: {e}")
            return file_path
    
    def _pptx_to_pdf(self, file_path: str, output_path: str) -> str:
        """PPTX转PDF（提取所有幻灯片为图片再合成PDF）"""
        try:
            from pptx import Presentation
            from PIL import Image
            import io
            
            prs = Presentation(file_path)
            
            # 获取幻灯片尺寸
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # 创建图片列表
            images = []
            for slide in prs.slides:
                # 创建空白图片
                img = Image.new('RGB', (int(slide_width/9525), int(slide_height/9525)), 'white')
                draw = ImageDraw.Draw(img)
                
                # 提取文字
                text_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text.strip())
                
                # 绘制文字
                y = 20
                for line in text_content[:20]:
                    if y > img.height - 40:
                        break
                    draw.text((20, y), line, fill='black')
                    y += 24
                
                images.append(img)
            
            if images:
                images[0].save(output_path, save_all=True, append_images=images[1:])
                logger.info(f"PPTX转PDF成功: {output_path}")
                return output_path
            
            return file_path
            
        except Exception as e:
            logger.error(f"PPTX转PDF失败: {e}")
            return file_path
    
    def _text_to_pdf(self, file_path: str, output_path: str) -> str:
        """文本转PDF"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            html_content = f"""
            <html>
            <body>
            <pre style="font-family: Arial, sans-serif; font-size: 12px;">
            {content}
            </pre>
            </body>
            </html>
            """
            
            from weasyprint import HTML
            HTML(string=html_content).write_pdf(output_path)
            return output_path
        
        except Exception as e:
            logger.error(f"文本转PDF失败: {e}")
            return file_path
    
    def delete_file(self, file_path: str) -> bool:
        """删除文件及其预览"""
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
                logger.info(f"已删除文件: {file_path}")
                
                # 同时删除预览文件
                preview_name = os.path.splitext(os.path.basename(file_path))[0]
                preview_path = os.path.join(self.preview_folder, f"{preview_name}.png")
                if os.path.exists(preview_path):
                    os.remove(preview_path)
                    logger.info(f"已删除预览: {preview_path}")
                
                return True
            return False
        except Exception as e:
            logger.error(f"删除文件失败: {e}")
            return False
    
    def list_files(self, folder: str = None) -> list:
        """列出文件夹中的文件"""
        folder = folder or self.upload_folder
        files = []
        
        try:
            for filename in os.listdir(folder):
                file_path = os.path.join(folder, filename)
                if os.path.isfile(file_path):
                    stat = os.stat(file_path)
                    files.append({
                        'filename': filename,
                        'path': file_path,
                        'size': stat.st_size,
                        'created': datetime.fromtimestamp(stat.st_ctime),
                        'modified': datetime.fromtimestamp(stat.st_mtime)
                    })
        
        except Exception as e:
            logger.error(f"列出文件失败: {e}")
        
        return sorted(files, key=lambda x: x['created'], reverse=True)
